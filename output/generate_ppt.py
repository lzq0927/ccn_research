#!/usr/bin/env python3
"""Generate 4+1 page PPT: Cross-industry reliability insights for cloud core network."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color palette ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_MID = RGBColor(0x16, 0x21, 0x3E)
BG_CARD = RGBColor(0x1E, 0x2D, 0x4A)
ACCENT_BLUE = RGBColor(0x00, 0x9C, 0xF0)
ACCENT_GREEN = RGBColor(0x00, 0xD6, 0x8E)
ACCENT_ORANGE = RGBColor(0xFF, 0x9F, 0x43)
ACCENT_RED = RGBColor(0xE8, 0x4D, 0x4D)
ACCENT_PURPLE = RGBColor(0xA2, 0x7B, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
TABLE_HEADER_BG = RGBColor(0x00, 0x5C, 0x8A)
TABLE_ROW_LIGHT = RGBColor(0x16, 0x21, 0x3E)
TABLE_ROW_DARK = RGBColor(0x12, 0x1C, 0x36)
INSIGHT_BG = RGBColor(0x0D, 0x3B, 0x5C)
TAG_SAME = RGBColor(0x00, 0x7A, 0x33)
TAG_PARTIAL = RGBColor(0xCC, 0x88, 0x00)
TAG_DIFF = RGBColor(0xCC, 0x33, 0x33)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_fill(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_text_box(slide, left, top, width, height, lines, default_size=10,
                      default_color=LIGHT_GRAY, line_spacing=1.3, font_name='Microsoft YaHei'):
    """lines: list of (text, font_size, color, bold) or just text string"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, fs, c, b = line, default_size, default_color, False
        else:
            text = line[0]
            fs = line[1] if len(line) > 1 else default_size
            c = line[2] if len(line) > 2 else default_color
            b = line[3] if len(line) > 3 else False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = font_name
        p.space_after = Pt(2)
        p.line_spacing = Pt(fs * line_spacing)
    return txBox


def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    """Add a styled table. headers: list of str. rows: list of list of str."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = 'Microsoft YaHei'
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for i, row in enumerate(rows):
        bg = TABLE_ROW_LIGHT if i % 2 == 0 else TABLE_ROW_DARK
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.color.rgb = LIGHT_GRAY
                p.font.name = 'Microsoft YaHei'
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape


def build_slide1(prs):
    """核电异构网"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)
    W = prs.slide_width
    H = prs.slide_height

    # Title bar
    add_shape_fill(slide, Inches(0), Inches(0), W, Inches(0.55), ACCENT_BLUE)
    add_text_box(slide, Inches(0.3), Inches(0.05), Inches(8), Inches(0.5),
                 "看行业：核电异构网", font_size=20, color=WHITE, bold=True)

    # ── Top: Requirements table ──
    add_text_box(slide, Inches(0.3), Inches(0.6), Inches(3), Inches(0.3),
                 "核电核心要求 vs 云核心网", font_size=11, color=ACCENT_BLUE, bold=True)

    headers = ["维度", "核电要求", "云核心网要求", "对比"]
    rows = [
        ["可用度", "CDF<10⁻⁵/堆年 (99.9999%)", "网元可用性≥99.999%", "相同(极致可靠)"],
        ["组网", "4重冗余+2oo4表决", "AMF Set/SMF Set多副本", "部分相同"],
        ["故障处置", "纵深防御5层+故障安全", "Pod重启→AZ容灾→Region容灾", "框架相同"],
        ["CCF防控", "异构冗余(能动+非能动)", "同构冗余为主", "不同(需补课)"],
        ["定期验证", "定期安全审查+PSA", "Chaos Engineering", "理念相同"],
    ]
    col_w = [Inches(0.9), Inches(3.0), Inches(2.8), Inches(1.6)]
    add_table(slide, Inches(0.3), Inches(0.9), Inches(8.3), Inches(2.0), headers, rows, col_w)

    # ── Middle: Left (tech) + Right (mapping) ──
    mid_top = Inches(3.05)
    mid_h = Inches(3.35)

    # Left panel - tech details
    add_shape_fill(slide, Inches(0.2), mid_top, Inches(5.2), Inches(0.28), ACCENT_BLUE)
    add_text_box(slide, Inches(0.3), mid_top, Inches(5), Inches(0.25),
                 "核电异构技术", font_size=10, color=WHITE, bold=True)

    left_lines = [
        ("1) 硬件异构", 10, ACCENT_GREEN, True),
        ("  策略A(最高): 数字化 vs 模拟/硬接线 — AP1000 PMS+DAS", 9, LIGHT_GRAY, False),
        ("  策略B(高):   不同CPU架构 CPU vs FPGA — NASPIC+NuPAC", 9, LIGHT_GRAY, False),
        ("  策略C(中):   不同制造商 — Common Q vs TXS vs Triconex", 9, LIGHT_GRAY, False),
        ("  华龙一号三重异构: FirmSys(CPU) + NASPIC(CPU) + NuPAC(FPGA)", 9, ACCENT_ORANGE, False),
        ("", 6, LIGHT_GRAY, False),
        ("2) 软件异构", 10, ACCENT_GREEN, True),
        ("  N版本编程(NVP): 多团队独立开发 → 表决逻辑", 9, LIGHT_GRAY, False),
        ("  Knight&Leveson: 85%概率无相似故障, 但不能保证独立", 9, LIGHT_GRAY, False),
        ("  最极端策略: AP1000 DAS — 非软件化模拟后备", 9, ACCENT_ORANGE, False),
        ("  形式化方法: B方法/SPARK Ada → 数学证明替代部分多样性", 9, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("3) 共因故障(CCF)量化与防控", 10, ACCENT_GREEN, True),
        ("  β因子模型: β=Q_CCF/Q_total → 多样性措施使β从0.3降至0.01", 9, LIGHT_GRAY, False),
        ("  策略A(不同技术): β降低80%-90%", 9, LIGHT_GRAY, False),
        ("  策略B(不同方法): β降低50%-70%", 9, LIGHT_GRAY, False),
        ("  ICDE数据库: 1000+CCF事件, 设计缺陷占30%", 9, LIGHT_GRAY, False),
        ("  关键洞察: 2oo3系统危险故障概率≈βλ, CCF主导!", 9, ACCENT_ORANGE, True),
    ]
    add_rich_text_box(slide, Inches(0.3), mid_top + Inches(0.3), Inches(5.0), Inches(3.0),
                      left_lines, default_size=9)

    # Right panel - cloud mapping
    add_shape_fill(slide, Inches(5.55), mid_top, Inches(3.25), Inches(0.28), ACCENT_PURPLE)
    add_text_box(slide, Inches(5.65), mid_top, Inches(3.0), Inches(0.25),
                 "云核心网借鉴", font_size=10, color=WHITE, bold=True)

    right_lines = [
        ("硬件异构 →", 10, ACCENT_PURPLE, True),
        ("  VM+K8s双栈部署(策略B)", 9, LIGHT_GRAY, False),
        ("  x86+ARM异构实例", 9, LIGHT_GRAY, False),
        ("  多云部署(策略A: β→0.01)", 9, LIGHT_GRAY, False),
        ("  4oo4表决 → Raft/Paxos共识", 9, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("软件异构 →", 10, ACCENT_PURPLE, True),
        ("  多版本容器镜像并行部署", 9, LIGHT_GRAY, False),
        ("  不同语言微服务(Go/Rust)", 9, LIGHT_GRAY, False),
        ("  不同基础镜像(Alpine/Ubuntu)", 9, LIGHT_GRAY, False),
        ("  协议Fuzzing替代部分多样性", 9, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("CCF防控 →", 10, ACCENT_PURPLE, True),
        ("  建立β因子基线评估", 9, LIGHT_GRAY, False),
        ("  渐进式: 同构→双栈→多云", 9, LIGHT_GRAY, False),
        ("  β=0.3→0.05→0.01 (3阶段)", 9, ACCENT_ORANGE, False),
        ("  可用性: 99.99%→99.995%→99.999%", 9, ACCENT_ORANGE, False),
    ]
    add_rich_text_box(slide, Inches(5.65), mid_top + Inches(0.3), Inches(3.1), Inches(3.0),
                      right_lines, default_size=9)

    # ── Bottom: Insight ──
    insight_top = Inches(6.55)
    add_shape_fill(slide, Inches(0.2), insight_top, Inches(8.6), Inches(0.75), INSIGHT_BG)
    add_text_box(slide, Inches(0.3), insight_top + Inches(0.02), Inches(1.5), Inches(0.25),
                 "洞察启示", font_size=10, color=ACCENT_ORANGE, bold=True)
    add_text_box(slide, Inches(0.3), insight_top + Inches(0.25), Inches(8.4), Inches(0.5),
                 "先防共因，再谈冗余 —— 2oo3系统危险故障概率≈βλ，CCF是可靠性的天花板 | "
                 "AI Coding大发展，高成本异构方案（N版本编程、多云部署）成为可能 | "
                 "核电纵深防御方法论 + 云Grid工程实践 = 核心网\"超可靠架构蓝图\"",
                 font_size=9, color=WHITE)


def build_slide2(prs):
    """云厂商分布式架构"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    W = prs.slide_width
    H = prs.slide_height

    add_shape_fill(slide, Inches(0), Inches(0), W, Inches(0.55), ACCENT_GREEN)
    add_text_box(slide, Inches(0.3), Inches(0.05), Inches(8), Inches(0.5),
                 "看行业：云厂商分布式架构", font_size=20, color=BG_DARK, bold=True)

    # Top table
    add_text_box(slide, Inches(0.3), Inches(0.6), Inches(3), Inches(0.3),
                 "云厂商可靠性要求 vs 云核心网", font_size=11, color=ACCENT_GREEN, bold=True)

    headers = ["维度", "云厂商实践", "云核心网需求", "对比"]
    rows = [
        ["可用度", "99.99%+(AWS/Azure SLA)", "99.999%(5个9目标)", "核心网更高"],
        ["故障隔离", "Cell-Based自治单元", "NF Set/切片隔离", "理念相同"],
        ["数据同步", "GaussDB RPO=0 RTO<10s", "UE/PDU Session状态同步", "部分相同"],
        ["切换速度", "秒级无感切换(UniformLive)", "AMF/UPF秒级切换", "目标相同"],
        ["弹性伸缩", "Auto Scaling+预测", "话务预测+NF弹性", "理念相同"],
    ]
    col_w = [Inches(0.9), Inches(3.0), Inches(2.8), Inches(1.6)]
    add_table(slide, Inches(0.3), Inches(0.9), Inches(8.3), Inches(2.0), headers, rows, col_w)

    # Middle
    mid_top = Inches(3.05)
    mid_h = Inches(3.35)

    add_shape_fill(slide, Inches(0.2), mid_top, Inches(5.2), Inches(0.28), ACCENT_GREEN)
    add_text_box(slide, Inches(0.3), mid_top, Inches(5), Inches(0.25),
                 "Cell化/Grid架构", font_size=10, color=BG_DARK, bold=True)

    left_lines = [
        ("1) AWS: 全球多Region+多AZ", 10, ACCENT_GREEN, True),
        ("  33个Region, 105个AZ, 600+边缘节点", 9, LIGHT_GRAY, False),
        ("  Aurora Global DB: 跨Region复制延迟<1s", 9, LIGHT_GRAY, False),
        ("  Route 53: 60s内检测端点故障", 9, LIGHT_GRAY, False),
        ("  Chaos Engineering验证弹性", 9, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("2) 华为云: 超级可用区+UniformLive", 10, ACCENT_GREEN, True),
        ("  超级AZ: 多DC融合, RTT<0.5ms, 故障透明切换", 9, LIGHT_GRAY, False),
        ("  GaussDB: RPO=0, RTO<10s, 邮储银行20亿笔/日", 9, LIGHT_GRAY, False),
        ("  UniformLive: 有状态应用跨Region多活", 9, ACCENT_ORANGE, False),
        ("  30分钟搭建Landing Zone", 9, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("3) Google: Spanner全球强一致", 10, ACCENT_GREEN, True),
        ("  TrueTime API: GPS+原子钟, 有界误差全局时钟", 9, LIGHT_GRAY, False),
        ("  Paxos: 3-5 Zone复制, 任意Zone故障不影响", 9, LIGHT_GRAY, False),
        ("  Anycast: 单IP全球可达, 亚秒级故障切换", 9, LIGHT_GRAY, False),
        ("  SRE方法论: 错误预算+SLO/SLI框架", 9, LIGHT_GRAY, False),
    ]
    add_rich_text_box(slide, Inches(0.3), mid_top + Inches(0.3), Inches(5.0), Inches(3.0),
                      left_lines, default_size=9)

    # Right
    add_shape_fill(slide, Inches(5.55), mid_top, Inches(3.25), Inches(0.28), ACCENT_PURPLE)
    add_text_box(slide, Inches(5.65), mid_top, Inches(3.0), Inches(0.25),
                 "云核心网借鉴", font_size=10, color=WHITE, bold=True)

    right_lines = [
        ("Cell化部署策略 →", 10, ACCENT_PURPLE, True),
        ("  三层: 中心Cell+区域Cell+边缘Cell", 9, LIGHT_GRAY, False),
        ("  渐进式: 同城Grid→跨Region→全国", 9, LIGHT_GRAY, False),
        ("  AMF Set映射到Grid Cell", 9, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("UPF跨Region无感切换 →", 10, ACCENT_PURPLE, True),
        ("  Anycast GTP-U: 单IP, BGP自动收敛", 9, LIGHT_GRAY, False),
        ("  会话状态同步: 增量同步<10ms", 9, LIGHT_GRAY, False),
        ("  SMF驱动: 500ms内完成切换", 9, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("数据一致性保障 →", 10, ACCENT_PURPLE, True),
        ("  UE/PDU Session: Paxos/Raft强一致", 9, LIGHT_GRAY, False),
        ("  计费数据: 最终一致性(分钟级)", 9, LIGHT_GRAY, False),
        ("  GaussDB方案: 分布式事务+跨AZ容灾", 9, ACCENT_ORANGE, False),
        ("", 6, LIGHT_GRAY, False),
        ("SRE体系借鉴 →", 10, ACCENT_PURPLE, True),
        ("  错误预算驱动交付节奏", 9, LIGHT_GRAY, False),
        ("  Chaos Engineering验证Grid有效性", 9, LIGHT_GRAY, False),
    ]
    add_rich_text_box(slide, Inches(5.65), mid_top + Inches(0.3), Inches(3.1), Inches(3.0),
                      right_lines, default_size=9)

    # Bottom insight
    insight_top = Inches(6.55)
    add_shape_fill(slide, Inches(0.2), insight_top, Inches(8.6), Inches(0.75), INSIGHT_BG)
    add_text_box(slide, Inches(0.3), insight_top + Inches(0.02), Inches(1.5), Inches(0.25),
                 "洞察启示", font_size=10, color=ACCENT_ORANGE, bold=True)
    add_text_box(slide, Inches(0.3), insight_top + Inches(0.25), Inches(8.4), Inches(0.5),
                 "单元化(Cell-Based)是故障爆炸半径控制的核心手段 —— 从应用级部署演进到基础设施级自治 | "
                 "数据面去中心化 + 控制面统一编排 = Grid架构的\"车之两轮\" | "
                 "有状态应用(AMF/SMF/UPF)的跨Region多活是核心网Grid化的关键技术难题",
                 font_size=9, color=WHITE)


def build_slide3(prs):
    """AI for 可靠性"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    W = prs.slide_width
    H = prs.slide_height

    add_shape_fill(slide, Inches(0), Inches(0), W, Inches(0.55), ACCENT_ORANGE)
    add_text_box(slide, Inches(0.3), Inches(0.05), Inches(8), Inches(0.5),
                 "看行业：AI for 可靠性", font_size=20, color=BG_DARK, bold=True)

    add_text_box(slide, Inches(0.3), Inches(0.6), Inches(3), Inches(0.3),
                 "AI可靠性能力 vs 云核心网需求", font_size=11, color=ACCENT_ORANGE, bold=True)

    headers = ["能力维度", "业界AI水平", "核心网适配", "对比"]
    rows = [
        ["异常检测", "Transformer/SSM秒级感知", "5GC信令KPI监测", "可直接迁移"],
        ["根因定位", "GNN+因果推断分钟级", "跨NF根因定位", "需适配拓扑"],
        ["故障预测", "时间序列预测(Transformer)", "话务预测+扩缩容", "可直接迁移"],
        ["自动自愈", "RL+Agent实验阶段", "NF自动重启/迁移", "需严格护栏"],
        ["LLM运维", "日志分析+修复建议", "3GPP规范+信令日志", "需领域微调"],
    ]
    col_w = [Inches(0.9), Inches(3.0), Inches(2.8), Inches(1.6)]
    add_table(slide, Inches(0.3), Inches(0.9), Inches(8.3), Inches(2.0), headers, rows, col_w)

    mid_top = Inches(3.05)
    mid_h = Inches(3.35)

    add_shape_fill(slide, Inches(0.2), mid_top, Inches(5.2), Inches(0.28), ACCENT_ORANGE)
    add_text_box(slide, Inches(0.3), mid_top, Inches(5), Inches(0.25),
                 "AI技术", font_size=10, color=BG_DARK, bold=True)

    left_lines = [
        ("1) AWS: 小模型→Agent化改造", 10, ACCENT_GREEN, True),
        ("  自动化推理: SMT求解器每天数十亿次推理检查", 9, LIGHT_GRAY, False),
        ("  小模型原子能力长期积累 → Agent快速编排适配", 9, LIGHT_GRAY, False),
        ("  DevOps Guru: ML驱动的云运维, 基于运营数据训练", 9, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("2) 微软: 仿真评估验证体系, 保障Agent输出可靠", 10, ACCENT_GREEN, True),
        ("  HALO: 层次化故障定位(服务→集群→实例)", 9, LIGHT_GRAY, False),
        ("  NENYA: 级联RL故障缓解, M365生产环境验证", 9, LIGHT_GRAY, False),
        ("  RESIN: 内存泄漏自动化处理(OSDI级别)", 9, LIGHT_GRAY, False),
        ("  Harness Engineering: Agent=Model+Harness管控", 9, ACCENT_ORANGE, False),
        ("  ACR权衡: Agency-Control-Reliability三角", 9, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("3) 终端云: 小模型→LLM/Agent→数据建模", 10, ACCENT_GREEN, True),
        ("  异常检测: iTransformer/DeMa/LEFT达到SOTA", 9, LIGHT_GRAY, False),
        ("  根因定位: Mulan多模态因果图(WWW'24, 被引97)", 9, LIGHT_GRAY, False),
        ("  LLM日志分析: R-Log推理增强, 提升228%", 9, LIGHT_GRAY, False),
        ("  Hermes Agent: 闭环自进化, 64K stars", 9, LIGHT_GRAY, False),
        ("  关键: 小模型/工具决定下限, Agent决定上限", 9, ACCENT_ORANGE, True),
    ]
    add_rich_text_box(slide, Inches(0.3), mid_top + Inches(0.3), Inches(5.0), Inches(3.0),
                      left_lines, default_size=9)

    add_shape_fill(slide, Inches(5.55), mid_top, Inches(3.25), Inches(0.28), ACCENT_PURPLE)
    add_text_box(slide, Inches(5.65), mid_top, Inches(3.0), Inches(0.25),
                 "云核心网借鉴", font_size=10, color=WHITE, bold=True)

    right_lines = [
        ("小模型原子能力 →", 10, ACCENT_PURPLE, True),
        ("  信令KPI异常检测(Isolation Forest+Transformer)", 9, LIGHT_GRAY, False),
        ("  告警聚合: 百万级→千级incident", 9, LIGHT_GRAY, False),
        ("  话务预测: PatchTST驱动预测性扩缩容", 9, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("Agent化改造 →", 10, ACCENT_PURPLE, True),
        ("  GNN根因定位: 5GC拓扑图+NRF服务注册", 9, LIGHT_GRAY, False),
        ("  LLM运维助手: 告警理解+3GPP规范查询", 9, LIGHT_GRAY, False),
        ("  自愈Agent: NF重启/迁移/降级决策", 9, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("Harness/仿真评估 →", 10, ACCENT_PURPLE, True),
        ("  三级置信度: 高→自动执行/中→人机协同/低→人工", 9, LIGHT_GRAY, False),
        ("  Chaos Engineering验证自愈有效性", 9, LIGHT_GRAY, False),
        ("  AIOpsLab基准测试评估算法选型", 9, LIGHT_GRAY, False),
        ("  Agent安全护栏: 权限白名单+操作审计", 9, ACCENT_ORANGE, False),
    ]
    add_rich_text_box(slide, Inches(5.65), mid_top + Inches(0.3), Inches(3.1), Inches(3.0),
                      right_lines, default_size=9)

    insight_top = Inches(6.55)
    add_shape_fill(slide, Inches(0.2), insight_top, Inches(8.6), Inches(0.75), INSIGHT_BG)
    add_text_box(slide, Inches(0.3), insight_top + Inches(0.02), Inches(1.5), Inches(0.25),
                 "洞察启示", font_size=10, color=ACCENT_ORANGE, bold=True)
    add_text_box(slide, Inches(0.3), insight_top + Inches(0.25), Inches(8.4), Inches(0.5),
                 "小模型/工具决定下限，Agent决定上限，Harness/仿真评估决定工程可用 | "
                 "AWS模式: 先积累小模型原子能力(自动化推理), 再Agent化编排 — 核心网应优先建设感知+诊断小模型 | "
                 "微软模式: 仿真评估体系是Agent可靠性的保障 — 核心网引入Agent必须有Chaos Engineering验证闭环",
                 font_size=9, color=WHITE)


def build_slide4(prs):
    """终端软件可靠性"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    W = prs.slide_width
    H = prs.slide_height

    add_shape_fill(slide, Inches(0), Inches(0), W, Inches(0.55), ACCENT_RED)
    add_text_box(slide, Inches(0.3), Inches(0.05), Inches(8), Inches(0.5),
                 "看行业：终端软件可靠性", font_size=20, color=WHITE, bold=True)

    add_text_box(slide, Inches(0.3), Inches(0.6), Inches(3), Inches(0.3),
                 "终端/软件可靠性 vs 云核心网", font_size=11, color=ACCENT_RED, bold=True)

    headers = ["维度", "终端/软件业界", "云核心网现状", "对比"]
    rows = [
        ["Bug发现", "静态分析+Fuzzing规模化(Google/Meta)", "协议测试为主", "需补课Fuzzing"],
        ["故障定界", "SBFL谱分析+LLM免测试定位", "人工日志分析", "需引入SBFL"],
        ["代码修复", "LLM Agent自动修复(RepairAgent)", "人工编码修复", "需引入APR"],
        ["持续保障", "Tricorder集成代码审查", "CI/CD基础流水线", "需深度集成"],
        ["分布式定位", "GNN+因果推断根因分析", "人工跨NF排查", "需引入AI"],
    ]
    col_w = [Inches(0.9), Inches(3.0), Inches(2.8), Inches(1.6)]
    add_table(slide, Inches(0.3), Inches(0.9), Inches(8.3), Inches(2.0), headers, rows, col_w)

    mid_top = Inches(3.05)
    mid_h = Inches(3.35)

    add_shape_fill(slide, Inches(0.2), mid_top, Inches(5.2), Inches(0.28), ACCENT_RED)
    add_text_box(slide, Inches(0.3), mid_top, Inches(5), Inches(0.25),
                 "故障定界定位修复技术", font_size=10, color=WHITE, bold=True)

    left_lines = [
        ("1) 苹果/谷歌: 问题定界→代码级定位→修复", 10, ACCENT_GREEN, True),
        ("  Google Tricorder: 静态分析嵌入代码审查, 90%+有用率", 9, LIGHT_GRAY, False),
        ("  Meta Infer: 分离逻辑+双推导, 数百万行/天", 9, LIGHT_GRAY, False),
        ("  LLM测试失败诊断: 已集成Google Critique审查系统", 9, LIGHT_GRAY, False),
        ("  RepairAgent: LLM Agent自主修复(ICSE'25, 被引398)", 9, ACCENT_ORANGE, False),
        ("", 6, LIGHT_GRAY, False),
        ("2) 鸿蒙/华为: 问题定界→代码级定位→修复", 10, ACCENT_GREEN, True),
        ("  DevEco Inspect: 静态分析工具集成", 9, LIGHT_GRAY, False),
        ("  分布式调试: 跨设备Trace + HiTrace链路追踪", 9, LIGHT_GRAY, False),
        ("  FaultLogger: 统一故障日志采集与聚合", 9, LIGHT_GRAY, False),
        ("  智能日志分析: NLP驱动的日志模式提取", 9, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("3) 5GC专项: 协议Fuzzing+一致性测试", 10, ACCENT_GREEN, True),
        ("  AFL-Net: 有状态网络协议灰盒Fuzzer", 9, LIGHT_GRAY, False),
        ("  5GC-Fuzz: 黑盒5GC协议测试(INFOCOM'25)", 9, LIGHT_GRAY, False),
        ("  CoreCrisis: 威胁引导5GC Fuzzing(USENIX'25)", 9, LIGHT_GRAY, False),
        ("  三级定位: 日志定界→SBFL定位→LLM诊断", 9, ACCENT_ORANGE, True),
    ]
    add_rich_text_box(slide, Inches(0.3), mid_top + Inches(0.3), Inches(5.0), Inches(3.0),
                      left_lines, default_size=9)

    add_shape_fill(slide, Inches(5.55), mid_top, Inches(3.25), Inches(0.28), ACCENT_PURPLE)
    add_text_box(slide, Inches(5.65), mid_top, Inches(3.0), Inches(0.25),
                 "云核心网借鉴", font_size=10, color=WHITE, bold=True)

    right_lines = [
        ("Bug发现体系 →", 10, ACCENT_PURPLE, True),
        ("  NGAP/PFCP/NAS协议Fuzzing进CI流水线", 9, LIGHT_GRAY, False),
        ("  SpotBugs/Clang静态分析嵌入代码审查", 9, LIGHT_GRAY, False),
        ("  参考Tricorder\"分析即服务\"模式", 9, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("故障定界定位 →", 10, ACCENT_PURPLE, True),
        ("  L1: OpenTelemetry+异常检测→网元定界", 9, LIGHT_GRAY, False),
        ("  L2: SBFL谱分析+LLM→代码行定位", 9, LIGHT_GRAY, False),
        ("  L3: RepairAgent→候选修复补丁", 9, LIGHT_GRAY, False),
        ("", 6, LIGHT_GRAY, False),
        ("持续可靠性保障 →", 10, ACCENT_PURPLE, True),
        ("  代码提交→静态分析+Fuzzing", 9, LIGHT_GRAY, False),
        ("  构建→增量分析+协议测试", 9, LIGHT_GRAY, False),
        ("  部署→灰度发布+端到端验证", 9, LIGHT_GRAY, False),
        ("  运维→AIOps智能运维+故障知识库", 9, LIGHT_GRAY, False),
    ]
    add_rich_text_box(slide, Inches(5.65), mid_top + Inches(0.3), Inches(3.1), Inches(3.0),
                      right_lines, default_size=9)

    insight_top = Inches(6.55)
    add_shape_fill(slide, Inches(0.2), insight_top, Inches(8.6), Inches(0.75), INSIGHT_BG)
    add_text_box(slide, Inches(0.3), insight_top + Inches(0.02), Inches(1.5), Inches(0.25),
                 "洞察启示", font_size=10, color=ACCENT_ORANGE, bold=True)
    add_text_box(slide, Inches(0.3), insight_top + Inches(0.25), Inches(8.4), Inches(0.5),
                 "问题定界 → 代码级定位 → 自修复 —— 三级流水线将故障处理从小时级压缩到分钟级 | "
                 "Google/Meta的\"分析即服务\"模式值得借鉴: 将静态分析无缝嵌入代码审查, 开发者零感知 | "
                 "5GC协议Fuzzing(NGAP/PFCP/NAS)是当前最大的可靠性盲区, 需优先建设",
                 font_size=9, color=WHITE)


def build_slide5_summary(prs):
    """Summary slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    W = prs.slide_width
    H = prs.slide_height

    # Full-width title
    add_shape_fill(slide, Inches(0), Inches(0), W, Inches(0.7), ACCENT_BLUE)
    add_text_box(slide, Inches(0.3), Inches(0.08), Inches(9), Inches(0.6),
                 "综合洞察：跨行业可靠性经验 → 云核心网技术演进路线", font_size=18, color=WHITE, bold=True)

    # Four insight cards
    card_w = Inches(4.15)
    card_h = Inches(2.6)
    gap = Inches(0.2)
    start_x = Inches(0.2)
    y1 = Inches(0.85)
    y2 = y1 + card_h + gap

    cards = [
        {
            "x": start_x, "y": y1,
            "color": ACCENT_BLUE, "title": "核电异构网",
            "lines": [
                ("核心洞察", 10, ACCENT_ORANGE, True),
                ("先防共因,再谈冗余 — β因子是可靠性天花板", 9, WHITE, False),
                ("关键借鉴", 10, ACCENT_ORANGE, True),
                ("纵深防御5层 → 核心网5层故障防线", 9, LIGHT_GRAY, False),
                ("异构冗余 → VM+K8s双栈→多云异构", 9, LIGHT_GRAY, False),
                ("β因子量化 → CCF防控效果可度量", 9, LIGHT_GRAY, False),
                ("能动+非能动 → 电讯级+IT级安全系统", 9, LIGHT_GRAY, False),
            ]
        },
        {
            "x": start_x + card_w + gap, "y": y1,
            "color": ACCENT_GREEN, "title": "云厂商分布式架构",
            "lines": [
                ("核心洞察", 10, ACCENT_ORANGE, True),
                ("单元化(Cell-Based)控制故障爆炸半径", 9, WHITE, False),
                ("关键借鉴", 10, ACCENT_ORANGE, True),
                ("超级AZ → AMF Set跨AZ部署", 9, LIGHT_GRAY, False),
                ("UniformLive → 有状态NF跨Region多活", 9, LIGHT_GRAY, False),
                ("GaussDB → Session数据分布式存储", 9, LIGHT_GRAY, False),
                ("SRE错误预算 → 可靠性量化管理", 9, LIGHT_GRAY, False),
            ]
        },
        {
            "x": start_x, "y": y2,
            "color": ACCENT_ORANGE, "title": "AI for 可靠性",
            "lines": [
                ("核心洞察", 10, ACCENT_ORANGE, True),
                ("小模型/工具决定下限, Agent决定上限", 9, WHITE, False),
                ("关键借鉴", 10, ACCENT_ORANGE, True),
                ("小模型原子能力 → 异常检测+根因定位", 9, LIGHT_GRAY, False),
                ("Agent化编排 → 自愈决策+修复执行", 9, LIGHT_GRAY, False),
                ("Harness/仿真评估 → Agent工程可用保障", 9, LIGHT_GRAY, False),
                ("三级置信度 → 高自动/中人机/低人工", 9, LIGHT_GRAY, False),
            ]
        },
        {
            "x": start_x + card_w + gap, "y": y2,
            "color": ACCENT_RED, "title": "终端软件可靠性",
            "lines": [
                ("核心洞察", 10, ACCENT_ORANGE, True),
                ("问题定界→代码定位→自修复 三级流水线", 9, WHITE, False),
                ("关键借鉴", 10, ACCENT_ORANGE, True),
                ("静态分析+协议Fuzzing → Bug发现", 9, LIGHT_GRAY, False),
                ("SBFL+LLM → 免测试代码级定位", 9, LIGHT_GRAY, False),
                ("RepairAgent → LLM自动程序修复", 9, LIGHT_GRAY, False),
                ("Tricorder模式 → 分析即服务嵌入审查", 9, LIGHT_GRAY, False),
            ]
        },
    ]

    for card in cards:
        add_shape_fill(slide, card["x"], card["y"], card_w, Inches(0.3), card["color"])
        add_text_box(slide, card["x"] + Inches(0.1), card["y"], card_w - Inches(0.2), Inches(0.28),
                     card["title"], font_size=11, color=WHITE if card["color"] in [ACCENT_BLUE, ACCENT_RED] else BG_DARK, bold=True)
        add_shape_fill(slide, card["x"], card["y"] + Inches(0.3), card_w, card_h - Inches(0.3), BG_CARD)
        add_rich_text_box(slide, card["x"] + Inches(0.15), card["y"] + Inches(0.35),
                          card_w - Inches(0.3), card_h - Inches(0.4),
                          card["lines"], default_size=9)

    # Bottom bar with timeline
    bar_top = y2 + card_h + Inches(0.15)
    add_shape_fill(slide, Inches(0.2), bar_top, Inches(8.6), Inches(0.85), INSIGHT_BG)
    add_text_box(slide, Inches(0.3), bar_top + Inches(0.02), Inches(2), Inches(0.25),
                 "演进路线", font_size=11, color=ACCENT_ORANGE, bold=True)
    add_text_box(slide, Inches(0.3), bar_top + Inches(0.28), Inches(8.4), Inches(0.55),
                 "短期(6月): 静态分析+异常检测+告警聚合+Fuzzing → 缺陷逃逸率降50% | "
                 "中期(1-2年): Grid同城多活+AI根因定位+LLM运维助手+VM/K8s双栈 → 可用性99.995% | "
                 "长期(3-5年): 全国Grid化+多云异构+AI全自动自愈+概率风险评估 → 可用性99.999%, β<0.01",
                 font_size=8, color=WHITE)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    build_slide1(prs)  # 核电异构网
    build_slide2(prs)  # 云厂商分布式架构
    build_slide3(prs)  # AI for 可靠性
    build_slide4(prs)  # 终端软件可靠性
    build_slide5_summary(prs)  # 综合洞察

    output_path = "D:/code/ccn_research/ccn_research/output/cross_industry_reliability_insights.pptx"
    prs.save(output_path)
    print(f"PPT saved to {output_path}")


if __name__ == "__main__":
    main()
