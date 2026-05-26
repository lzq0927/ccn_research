#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
云核心网可靠性技术洞察 PPT v2
白色简洁背景 · 少量配色 · 引用源文架构图
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── 配色（简洁双色系 + 灰度）──────────────────────────
PRIMARY    = RGBColor(0x1A, 0x56, 0xDB)   # 主色：深蓝
SECONDARY  = RGBColor(0x37, 0x41, 0x51)   # 辅色：深灰
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BG_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xF3, 0xF4, 0xF6)   # 浅灰背景块
BORDER     = RGBColor(0xD1, 0xD5, 0xDB)   # 边框灰
DARK_TEXT   = RGBColor(0x11, 0x18, 0x27)   # 正文黑
MID_TEXT    = RGBColor(0x6B, 0x72, 0x80)   # 次要文字灰
ACCENT_RED = RGBColor(0xDC, 0x26, 0x26)   # 警示红
ACCENT_GRN = RGBColor(0x16, 0xA3, 0x4A)   # 成功绿
SOURCE_CLR = RGBColor(0x9C, 0xA3, 0xAF)   # 来源说明

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── 工具函数 ──────────────────────────────────────────

def set_bg(slide, color=BG_WHITE):
    bg = slide.background; f = bg.fill; f.solid(); f.fore_color.rgb = color

def tb(slide, l, t, w, h, text, sz=14, clr=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, font='Microsoft YaHei'):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = clr; p.font.bold = bold
    p.font.name = font; p.alignment = align
    return box

def rich_tb(slide, l, t, w, h):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    return tf

def para(tf, text, sz=12, clr=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, sb=Pt(0), sa=Pt(4), font='Microsoft YaHei'):
    p = tf.paragraphs[0] if len(tf.paragraphs)==1 and tf.paragraphs[0].text=='' else tf.add_paragraph()
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = clr
    p.font.bold = bold; p.font.name = font; p.alignment = align
    p.space_before = sb; p.space_after = sa
    return p

def rect(slide, l, t, w, h, fill=LIGHT_BG, border=BORDER):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(0.75)
    s.adjustments[0] = 0.03
    return s

def rect_text(s, text, sz=11, clr=DARK_TEXT, bold=False, align=PP_ALIGN.CENTER):
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = clr; p.font.bold = bold
    p.font.name = 'Microsoft YaHei'; p.alignment = align

def arrow(slide, l, t, w, h, clr=PRIMARY):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()
    return s

def chevron(slide, l, t, w, h, clr=PRIMARY):
    s = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()
    return s

def src_note(slide, text, t=Inches(6.95)):
    tb(slide, Inches(0.5), t, Inches(12.3), Inches(0.4),
       f"[图源] {text}", sz=8, clr=SOURCE_CLR)

def header(slide, title, sub=""):
    # 顶部蓝色条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.88))
    bar.fill.solid(); bar.fill.fore_color.rgb = PRIMARY; bar.line.fill.background()
    tb(slide, Inches(0.6), Inches(0.12), Inches(12), Inches(0.55),
       title, sz=22, clr=WHITE, bold=True)
    if sub:
        tb(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.32),
           sub, sz=11, clr=RGBColor(0xBF,0xDB,0xFE))

def pg(slide, n, total=18):
    tb(slide, Inches(12.3), Inches(7.1), Inches(0.8), Inches(0.3),
       f"{n}/{total}", sz=8, clr=SOURCE_CLR, align=PP_ALIGN.RIGHT)

# ── 第1页：封面 ───────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
# 大色块背景
bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), SLIDE_W, Inches(2.8))
bar.fill.solid(); bar.fill.fore_color.rgb = PRIMARY; bar.line.fill.background()
tb(sl, Inches(1.2), Inches(1.2), Inches(11), Inches(0.5),
   "云核心网可靠性技术洞察", sz=38, clr=DARK_TEXT, bold=True)
tb(sl, Inches(1.2), Inches(1.8), Inches(11), Inches(0.5),
   "跨界借鉴与前沿技术追踪", sz=18, clr=PRIMARY)
tb(sl, Inches(1.2), Inches(3.0), Inches(11), Inches(0.8),
   "核电异构架构 · 云Grid架构 · AI for可靠性 · 软件可靠性", sz=20, clr=WHITE, bold=True)
tb(sl, Inches(1.2), Inches(3.8), Inches(11), Inches(0.5),
   "基于CCF A/B类学术文献 · IEC/IAEA/3GPP行业标准 · 头部厂商技术报告", sz=12, clr=RGBColor(0xBF,0xDB,0xFE))
tb(sl, Inches(1.2), Inches(5.8), Inches(5), Inches(0.4),
   "2026年5月", sz=14, clr=MID_TEXT)
tb(sl, Inches(8), Inches(5.8), Inches(4.5), Inches(0.4),
   "4份独立报告 + 1份综合分析  |  93+引用文献  |  195KB+深度内容", sz=11, clr=MID_TEXT, align=PP_ALIGN.RIGHT)
pg(sl, 1)

# ── 第2页：执行摘要 ────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header(sl, "执行摘要", "面向高层决策者的四个方向核心发现")

findings = [
    ("Track 1: 核电异构架构",
     "• 纵深防御五层架构(IAEA SSR-2/1)\n• Beta因子量化: β 0.3→0.01(97%↓)\n• 堆芯损坏频率CDF < 10⁻⁷/堆年\n• \"能动+非能动\"→VM+K8s双栈容灾"),
    ("Track 2: 云Grid架构",
     "• 华为云Super AZ/UniformLive秒级切换\n• Cell-Based自治→5GC三层Cell部署\n• AMF Set跨Region + Anycast GTP-U\n• 数据一致性: Paxos/Raft强一致"),
    ("Track 3: AI for 可靠性",
     "• 张圣林/裴丹团队最新成果(TOSEM/KDD)\n• 微软MSRA 25篇AIOps论文(HALO/NENYA)\n• Hermes Agent自进化+Harness Engineering\n• LLM运维(R-Log/Xpert/UniLog)"),
    ("Track 4: 软件可靠性",
     "• GNN/因果推断根因定位 + LLM免测试FL\n• RepairAgent(ICSE'25)自动程序修复\n• 5GC-Fuzz(CoreCrisis)专项协议Fuzzing\n• Google Tricorder/Meta Infer工业实践"),
]
for i, (title, desc) in enumerate(findings):
    x = Inches(0.4 + i * 3.2)
    card = rect(sl, x, Inches(1.1), Inches(3.0), Inches(4.0), fill=WHITE)
    # 标题条
    tbar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.1), Inches(3.0), Inches(0.45))
    tbar.fill.solid(); tbar.fill.fore_color.rgb = PRIMARY; tbar.line.fill.background()
    rect_text(tbar, title, sz=12, clr=WHITE, bold=True)
    tb(sl, x+Inches(0.15), Inches(1.65), Inches(2.7), Inches(3.3),
       desc, sz=11, clr=DARK_TEXT)

# 底部量化目标
tb(sl, Inches(0.5), Inches(5.4), Inches(12), Inches(0.3),
   "量化目标", sz=14, clr=PRIMARY, bold=True)
targets = ["可用性 99.999%", "MTTD < 30秒", "MTTR < 5分钟", "Beta因子 < 0.01", "AI定位率 >95%"]
for i, t in enumerate(targets):
    r = rect(sl, Inches(0.5+i*2.5), Inches(5.8), Inches(2.3), Inches(0.5), fill=PRIMARY, border=PRIMARY)
    rect_text(r, t, sz=12, clr=WHITE, bold=True)
pg(sl, 2)

# ── 第3页：Track 1 核电异构架构 ───────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header(sl, "Track 1: 核电异构架构洞察", "SIL 4级 · 纵深防御 · 异构冗余 · 共因故障防控")

# 左：纵深防御
tb(sl, Inches(0.5), Inches(1.1), Inches(6), Inches(0.3),
   "纵深防御五层架构 (Defense-in-Depth)", sz=14, clr=PRIMARY, bold=True)
did = [
    ("L5 全局应急", "场外应急响应、公众防护措施"),
    ("L4 严重事故缓解", "严重事故管理导则、附加安全设施"),
    ("L3 设计基准事故控制", "反应堆保护系统、工程安全设施"),
    ("L2 偏差控制", "限制系统、自动控制系统"),
    ("L1 预防", "高质量设计建造、保守设计裕量"),
]
colors = [RGBColor(0xDC,0x26,0x26), RGBColor(0xEA,0x58,0x3C), RGBColor(0xD9,0x77,0x06),
          RGBColor(0x16,0xA3,0x4A), RGBColor(0x15,0x80,0x3D)]
for i, (name, desc) in enumerate(did):
    y = Inches(1.5 + i*0.72)
    r = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(5.5-i*0.25), Inches(0.58))
    r.fill.solid(); r.fill.fore_color.rgb = colors[i]; r.line.fill.background()
    rect_text(r, f"{name}  —  {desc}", sz=11, clr=WHITE)

# 右：Beta因子递减
tb(sl, Inches(6.8), Inches(1.1), Inches(6), Inches(0.3),
   "Beta因子递减路径（共因故障量化）", sz=14, clr=PRIMARY, bold=True)
steps = [
    ("未防控", "β≈0.30", "同构冗余，CCF概率高"),
    ("镜像多样性", "β≈0.15", "多版本镜像并行，↓50%"),
    ("技术栈异构", "β≈0.05", "VM+K8s双栈，↓83%"),
    ("多云部署", "β≈0.01", "跨云/混合云，↓97%"),
]
sc = [RGBColor(0xF3,0xF4,0xF6), RGBColor(0xDB,0xED,0xE4), RGBColor(0xA7,0xF3,0xD0), RGBColor(0x6E,0xE7,0xB7)]
tc = [ACCENT_RED, RGBColor(0xD9,0x77,0x06), ACCENT_GRN, RGBColor(0x04,0x78,0x57)]
for i, (stage, beta, desc) in enumerate(steps):
    y = Inches(1.5 + i*1.05)
    r = rect(sl, Inches(6.8), y, Inches(5.8), Inches(0.85), fill=sc[i])
    tf = rich_tb(sl, Inches(7.0), y+Inches(0.05), Inches(5.4), Inches(0.75))
    para(tf, f"{stage}  {beta}", sz=13, clr=tc[i], bold=True)
    para(tf, desc, sz=10, clr=MID_TEXT, sa=Pt(0))

src_note(sl, "IAEA SSR-2/1(2016) Fig.1 Defence-in-depth; NUREG/CR-6303(2009) Fig.3-1 Diversity strategies; "
    "IEC 61508:2010 Part6 表B/C 共因故障Beta因子; IAEA Safety Reports No.90 Fig.5 CCF防御策略")
pg(sl, 3)

# ── 第4页：Track 2 云Grid架构 ─────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header(sl, "Track 2: 云Grid架构洞察", "Cell-Based自治 · 华为云Super AZ · 秒级无感切换")

# 演进路线
tb(sl, Inches(0.5), Inches(1.1), Inches(12), Inches(0.3),
   "架构演进路线", sz=14, clr=PRIMARY, bold=True)
evo = [("主备模式","分钟~小时RTO"),("同城双活","秒级RTO"),("异地多活","分钟RTO"),("Grid架构","秒级/无感")]
for i, (n, d) in enumerate(evo):
    x = Inches(0.5+i*3.15)
    if i > 0: chevron(sl, x-Inches(0.25), Inches(1.55), Inches(0.4), Inches(0.35), BORDER)
    r = rect(sl, x, Inches(1.45), Inches(2.9), Inches(0.6))
    rect_text(r, n, sz=14, clr=PRIMARY, bold=True)
    tb(sl, x, Inches(2.1), Inches(2.9), Inches(0.3), d, sz=9, clr=MID_TEXT, align=PP_ALIGN.CENTER)

# 华为云Grid
tb(sl, Inches(0.5), Inches(2.5), Inches(6), Inches(0.3),
   "华为云Grid架构层次", sz=14, clr=PRIMARY, bold=True)
hw = [
    ("全局控制面", "RGC(Resource Governance Center) — 统一资源编排"),
    ("Grid调度层", "智能调度引擎 — 跨Cell负载均衡、故障转移"),
    ("Cell自治层", "Super AZ — 自包含计算/存储/网络"),
    ("数据同步层", "DRS + GaussDB分布式 — 跨Cell强一致性"),
    ("服务治理层", "CSE微服务引擎 — 服务注册、熔断、灰度发布"),
]
for i, (n, d) in enumerate(hw):
    y = Inches(2.9+i*0.65)
    r = rect(sl, Inches(0.5), y, Inches(5.8), Inches(0.55))
    tf = rich_tb(sl, Inches(0.65), y+Inches(0.02), Inches(5.5), Inches(0.5))
    para(tf, n, sz=11, clr=PRIMARY, bold=True, sa=Pt(0))
    para(tf, d, sz=9, clr=MID_TEXT, sa=Pt(0))

# 厂商对比
tb(sl, Inches(6.8), Inches(2.5), Inches(6), Inches(0.3),
   "厂商对比矩阵", sz=14, clr=PRIMARY, bold=True)
vendors = [
    ("华为云", "Super AZ, UniformLive\nGaussDB, CSE, DRS", "电信级+互联网弹性"),
    ("AWS", "Multi-AZ, Aurora Global\nDevOps Guru, Auto Reasoning", "最成熟云生态"),
    ("Google", "Borg/Omega, Spanner\nSRE体系, ACA", "全球分布式先驱"),
    ("Azure", "AZ, Cosmos DB\nChaos Studio", "企业级混合云"),
]
for i, (n, t, d) in enumerate(vendors):
    col = i%2; row = i//2
    x = Inches(6.8+col*3.1); y = Inches(2.9+row*1.4)
    r = rect(sl, x, y, Inches(2.9), Inches(1.2))
    tf = rich_tb(sl, x+Inches(0.1), y+Inches(0.05), Inches(2.7), Inches(1.1))
    para(tf, n, sz=13, clr=PRIMARY, bold=True, sa=Pt(2))
    para(tf, t, sz=9, clr=DARK_TEXT, sa=Pt(2))
    para(tf, d, sz=9, clr=MID_TEXT, sa=Pt(0))

src_note(sl, "华为云分布式云战略白皮书(2023) Fig.2 Grid架构; Verma et al., Borg, EuroSys 2015 Fig.1; "
    "Corbett et al., Spanner, TOCS 2013 Fig.1; AWS Well-Architected Reliability Pillar(2024)")
pg(sl, 4)

# ── 第5页：Track 2 Detail — 核心网Grid化 ──────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header(sl, "5G核心网Grid化部署方案", "三层Cell架构 · UPF无感切换 · AMF/SMF异地多活")

# 三层Cell
cells = [
    ("中心Cell", "UDM · AUSF · NRF\n全国/省级控制面\n核心DC", Inches(1.2)),
    ("区域Cell", "AMF · SMF · UPF · PCF\n区域控制面+用户面\n地市DC", Inches(3.0)),
    ("边缘Cell", "边缘UPF · MEC应用\n超低延迟处理\n基站侧/边缘DC", Inches(4.8)),
]
for n, d, y in cells:
    r = rect(sl, Inches(1.5), y, Inches(4.5), Inches(1.4))
    rect_text(r, f"{n}\n{d}", sz=12, clr=DARK_TEXT)

# 右侧关键方案
tb(sl, Inches(7), Inches(1.1), Inches(5.5), Inches(0.3),
   "关键技术方案", sz=14, clr=PRIMARY, bold=True)
schemes = [
    ("UPF跨Region无感切换", "Anycast GTP-U + 会话状态同步\nRPO=0, RTO<1s", "参考: 华为云UniformLive架构"),
    ("AMF异地多活", "AMF Set跨Region + SCTP Multi-homing\nGUAMI动态重分配", "参考: 3GPP TS 23.501 §5.5"),
    ("Session数据一致性", "强一致: Paxos/Raft (UE上下文)\n最终一致: 异步复制 (计费数据)", "参考: GaussDB分布式事务"),
]
for i, (n, d, s) in enumerate(schemes):
    y = Inches(1.5+i*1.5)
    r = rect(sl, Inches(7), y, Inches(5.5), Inches(1.3))
    tf = rich_tb(sl, Inches(7.15), y+Inches(0.05), Inches(5.2), Inches(1.2))
    para(tf, n, sz=12, clr=PRIMARY, bold=True, sa=Pt(2))
    para(tf, d, sz=10, clr=DARK_TEXT, sa=Pt(2))
    para(tf, s, sz=8, clr=SOURCE_CLR, sa=Pt(0))

src_note(sl, "3GPP TS 23.501 v18.4.0 Fig.4.2.3-1 5GC架构; TS 23.502 Fig.4.2.2 Registration; "
    "华为云DRS数据复制服务 huaweicloud.com/product/drs")
pg(sl, 5)

# ── 第6页：Track 3 AI for 可靠性 ──────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header(sl, "Track 3: AI for 可靠性洞察", "感知-诊断-预测-自愈闭环 · 张圣林/裴丹/微软MSRA最新成果")

# AI闭环
loop = [
    ("感知", "异常检测\n告警聚合", "iTransformer\nPaAno(ICLR'26)\nAnomaly Transformer", ACCENT_GRN),
    ("诊断", "根因定位\n故障溯源", "MicroRCA(SIGCOMM)\nMulan(WWW'24)\nTraceRCA(IWQOS)", PRIMARY),
    ("预测", "故障预警\n预测扩缩容", "PatchTST\nTimesFM\nCAD(KDD'23)", RGBColor(0xD9,0x77,0x06)),
    ("自愈", "自动恢复\n容灾决策", "NENYA(KDD'22)\nHermes Agent\nRL self-healing", ACCENT_RED),
]
for i, (n, d, t, c) in enumerate(loop):
    x = Inches(0.5+i*3.2)
    # 卡片
    r = rect(sl, x, Inches(1.1), Inches(2.95), Inches(3.0))
    # 标题圈
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, x+Inches(0.75), Inches(1.2), Inches(1.3), Inches(1.3))
    circ.fill.solid(); circ.fill.fore_color.rgb = c; circ.line.fill.background()
    rect_text(circ, n, sz=20, clr=WHITE, bold=True)
    tb(sl, x+Inches(0.15), Inches(2.6), Inches(2.65), Inches(0.5), d, sz=11, clr=DARK_TEXT, align=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.15), Inches(3.1), Inches(2.65), Inches(0.8), t, sz=9, clr=PRIMARY, align=PP_ALIGN.CENTER)
    if i < 3:
        arrow(sl, x+Inches(2.95), Inches(1.6), Inches(0.35), Inches(0.3), BORDER)

# 底部：权威学者
tb(sl, Inches(0.5), Inches(4.3), Inches(12), Inches(0.3),
   "权威学术研究（2024-2026高影响力论文）", sz=13, clr=PRIMARY, bold=True)

scholars = [
    ("张圣林(南开)", "多模态诊断(TSC'23,125引) · 时序知识图谱(TSC'24)\nAIOpsLab(FSE'25) · R-Log LLM日志(arXiv'25)\n微服务故障综述(TOSEM'25,85引)", "图源: Zhang et al., TOSEM 2025 Fig.1 诊断方法分类; TSC 2024 Fig.2 知识图谱架构"),
    ("裴丹(清华)", "KPI预训练(KDD'24) · TraceRCA(IWQOS'21,223引)\n多模态自适应(ICSE'24,31引) · TimeSeriesBench(ISSRE'24)\n日志解析基准(ICSE'19,706引)", "图源: Ma et al., ICSE 2024 Fig.3 多模态框架; KDD 2024 Fig.2 预训练架构"),
    ("微软MSRA", "HALO层次定位(KDD'21) · NENYA级联RL(KDD'22)\nRESIN内存泄漏(OSDI'22) · Xpert LLM事件(ICSE'24)\nSPINE日志解析(FSE'22 Best Paper)", "图源: NENYA, KDD 2022 Fig.2 级联RL架构; HALO, KDD 2021 Fig.3 层次定位"),
]
for i, (n, d, s) in enumerate(scholars):
    y = Inches(4.65+i*0.75)
    r = rect(sl, Inches(0.5), y, Inches(12.3), Inches(0.65))
    tf = rich_tb(sl, Inches(0.6), y+Inches(0.02), Inches(2.0), Inches(0.6))
    para(tf, n, sz=10, clr=PRIMARY, bold=True, sa=Pt(0))
    tf2 = rich_tb(sl, Inches(2.6), y+Inches(0.02), Inches(6.5), Inches(0.6))
    para(tf2, d, sz=8, clr=DARK_TEXT, sa=Pt(0))
    tb(sl, Inches(9.2), y+Inches(0.02), Inches(3.5), Inches(0.6), s, sz=7, clr=SOURCE_CLR)

pg(sl, 6)

# ── 第7页：Track 3 Detail — AI Agent工程 ──────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header(sl, "AI Agent工程实践与可靠性保障", "Hermes Agent · Harness Engineering · Claude Code · OpenClaw")

agents = [
    ("Hermes Agent\n(Nous Research, 2026.2)",
     "• 64K+ GitHub stars, 开源自进化AI Agent\n• 闭环学习: DSPy+GEPA自动进化技能\n• 四层记忆系统, 跨会话持久化\n• 3个CVE(2026.4)警示安全边界",
     "图源: github.com/nousresearch/hermes-agent README架构图;\n"
     "mranand.substack.com 技术深文 Fig.1 Learning Loop; Fig.2 Memory Layers"),
    ("Harness Engineering\n(2025-2026新兴学科)",
     "• Agent = Model + Harness 理念\n• arXiv 2602.16666: 12指标Agent可靠性\n• CAR框架: Control·Agency·Runtime\n• OpenAI Harness: 仓库脚手架+CI管控",
     "图源: arxiv 2602.16666v2 Fig.2 可靠性指标分解;\n"
     "preprints.org 202603.1756 Fig.1 CAR分解; openai.com/index/harness-engineering"),
    ("Claude Code\n(Anthropic, #1 Agent)",
     "• arXiv 2604.14228 设计空间论文\n• 测试驱动: 实现→运行测试→调试→迭代\n• MCP协议: 连接Sentry/Datadog/K8s\n• 管道式日志: tail -f | claude -p",
     "图源: arxiv 2604.14228v1 Fig.2 Agent设计空间;\n"
     "docs.anthropic.com/claude-code 架构图; MCP protocol modelcontextprotocol.io"),
    ("OpenClaw\n(150K-280K stars, 2025.11)",
     "• 本地网关架构, 自托管运行\n• WhatsApp/Telegram/Slack集成\n• Active Memory + Task Brain(2026)\n• 创建者加入OpenAI后转独立治理",
     "图源: openclaw.ai 官网架构图;\n"
     "tosea.ai/blog/openclaw Fig.1 Local Gateway架构"),
]
for i, (n, d, s) in enumerate(agents):
    col = i%2; row = i//2
    x = Inches(0.5+col*6.4); y = Inches(1.1+row*3.0)
    r = rect(sl, x, y, Inches(6.1), Inches(2.7))
    tf = rich_tb(sl, x+Inches(0.15), y+Inches(0.05), Inches(5.8), Inches(2.6))
    para(tf, n, sz=13, clr=PRIMARY, bold=True, sa=Pt(4))
    para(tf, d, sz=10, clr=DARK_TEXT, sa=Pt(6))
    para(tf, s, sz=7, clr=SOURCE_CLR, sa=Pt(0))

pg(sl, 7)

# ── 第8页：Track 4 软件可靠性 ────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header(sl, "Track 4: 软件可靠性洞察", "Bug发现·定界定位·自动修复 — 全链路软件可靠性工程")

# 流水线
tb(sl, Inches(0.5), Inches(1.1), Inches(12), Inches(0.3),
   "全链路Bug管理流水线", sz=14, clr=PRIMARY, bold=True)
pipe = [
    ("发现", "静态分析\nFuzzing", "Tricorder\nAFL-Net\n5GC-Fuzz"),
    ("定位", "SBFL\nLLM定位", "Ochiai\nAgentFL\nRepairAgent"),
    ("修复", "APR\nLLM Agent", "ICSE'25\nDebugRepair"),
    ("验证", "回归测试\n灰度发布", "Harness CV\n金丝雀部署"),
]
for i, (n, d, t) in enumerate(pipe):
    x = Inches(0.5+i*3.2)
    r = rect(sl, x, Inches(1.5), Inches(2.95), Inches(1.6))
    tb(sl, x+Inches(0.1), Inches(1.55), Inches(2.75), Inches(0.35),
       n, sz=16, clr=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.1), Inches(1.95), Inches(2.75), Inches(0.5),
       d, sz=11, clr=DARK_TEXT, align=PP_ALIGN.CENTER)
    tb(sl, x+Inches(0.1), Inches(2.5), Inches(2.75), Inches(0.5),
       t, sz=9, clr=MID_TEXT, align=PP_ALIGN.CENTER)
    if i < 3: arrow(sl, x+Inches(2.95), Inches(1.95), Inches(0.35), Inches(0.3), BORDER)

# 工业实践
tb(sl, Inches(0.5), Inches(3.3), Inches(12), Inches(0.3),
   "工业界大规模实践", sz=13, clr=PRIMARY, bold=True)
practices = [
    ("Google Tricorder", "数百种静态分析器嵌入代码审查\n每日分析数十亿行代码\n2026: LLM集成测试诊断(Critique)\n图源: Sadowski, CACM 2018 Fig.1"),
    ("Meta Infer", "分离逻辑+双推导分析\n每日扫描数百万行代码\nESOP 2024增强组合式框架\n图源: Distefano, CACM 2019 Fig.2"),
    ("5GC专项Fuzzing", "5GC-Fuzz(INFOCOM'25)黑盒测试\nCoreCrisis(USENIX Sec'25)\nAFL-Net有状态协议Fuzzing\n图源: Sun, INFOCOM 2025 Fig.2 5GC测试架构"),
]
for i, (n, d) in enumerate(practices):
    x = Inches(0.5+i*4.15)
    r = rect(sl, x, Inches(3.7), Inches(3.9), Inches(1.8))
    tf = rich_tb(sl, x+Inches(0.1), Inches(3.75), Inches(3.7), Inches(1.7))
    para(tf, n, sz=13, clr=PRIMARY, bold=True, sa=Pt(3))
    para(tf, d, sz=9, clr=DARK_TEXT, sa=Pt(0))

# 底部：核心网DevOps
r = rect(sl, Inches(0.4), Inches(5.7), Inches(12.5), Inches(0.9))
rect_text(r, "核心网DevOps集成: 代码提交(静态分析) → 构建(增量Fuzzing) → 测试(一致性+接口Fuzzing) → 部署(灰度+智能回滚) → 运维(AIOps根因定位)",
          sz=11, clr=PRIMARY, bold=True)
pg(sl, 8)

# ── 第9页：综合架构蓝图 ──────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header(sl, "综合洞察：核心网超可靠架构蓝图", "核电纵深防御 × 云Grid工程 × AI全链路赋能")

layers = [
    ("第5层 全局应急", "DNS全局调度 · 灾难恢复 · 事后分析", "RTO<30min",
     "图源: IAEA SSR-2/1 DiD第5层; 华为云GSLB架构"),
    ("第4层 Region容灾", "跨Region切换 · 异构冗余(β<0.01) · 数据恢复", "RTO<5min",
     "图源: 华为云Grid Cell架构 Fig.3; NENYA KDD'22 RL容灾"),
    ("第3层 AZ容灾", "Pod自动重启 · AZ内切换 · 熔断降级", "RTO<1min",
     "图源: Kubernetes Pod生命周期; Netflix Chaos Monkey实验"),
    ("第2层 检测控制", "AI异常检测(Transformer) · 限流降级 · 多层健康检查", "MTTD<30s",
     "图源: Anomaly Transformer ICLR'22 Fig.1; iTransformer Fig.2"),
    ("第1层 代码预防", "静态分析+Fuzzing · 代码审查 · CI/CD流水线", "缺陷↓50%",
     "图源: Sadowski, CACM 2018 Fig.1 Tricorder架构"),
]
colors = [RGBColor(0xDC,0x26,0x26), RGBColor(0xEA,0x58,0x3C), RGBColor(0xD9,0x77,0x06),
          RGBColor(0x16,0xA3,0x4A), RGBColor(0x15,0x80,0x3D)]
for i, (n, d, tgt, src) in enumerate(layers):
    y = Inches(1.1+i*0.98)
    w = Inches(8.5-i*0.3)
    r = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, w, Inches(0.82))
    r.fill.solid(); r.fill.fore_color.rgb = colors[i]; r.line.fill.background()
    rect_text(r, f"{n}  —  {d}", sz=11, clr=WHITE)
    # 目标
    tr = rect(sl, Inches(0.5)+w+Inches(0.2), y+Inches(0.12), Inches(1.4), Inches(0.5), fill=LIGHT_BG)
    rect_text(tr, tgt, sz=10, clr=colors[i], bold=True)
    # 来源
    tb(sl, Inches(10.5), y+Inches(0.1), Inches(2.5), Inches(0.65), src, sz=7, clr=SOURCE_CLR)

pg(sl, 9)

# ── 第10页：技术融合闭环 ─────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header(sl, "综合洞察：开发态 × 运行态 全链路闭环", "AI for可靠性 + 软件可靠性 → 分钟级故障修复全流程")

flow = [
    ("运行时异常检测", "Track 3\nTransformer/SSM", ACCENT_GRN),
    ("根因网元定位", "Track 3\nGNN/因果推断", PRIMARY),
    ("代码级故障定位", "Track 4\nSBFL/LLM", RGBColor(0xD9,0x77,0x06)),
    ("自动代码修复", "Track 4\nRepairAgent", ACCENT_RED),
    ("验证部署上线", "Track 4+3\nCI/CD+Harness", RGBColor(0x7C,0x3A,0xED)),
]
for i, (n, t, c) in enumerate(flow):
    x = Inches(0.3+i*2.6)
    r = rect(sl, x, Inches(1.2), Inches(2.35), Inches(1.4), fill=LIGHT_BG)
    rect_text(r, f"{n}\n\n{t}", sz=10, clr=DARK_TEXT)
    if i < 4: arrow(sl, x+Inches(2.35), Inches(1.6), Inches(0.35), Inches(0.3), BORDER)

# 两种模式对比
tb(sl, Inches(0.5), Inches(2.9), Inches(12), Inches(0.3),
   "运维范式转变", sz=14, clr=PRIMARY, bold=True)

r1 = rect(sl, Inches(0.5), Inches(3.3), Inches(5.8), Inches(1.6), fill=RGBColor(0xFE,0xF2,0xF2))
tf1 = rich_tb(sl, Inches(0.65), Inches(3.35), Inches(5.5), Inches(1.5))
para(tf1, "传统模式 (被动响应)", sz=13, clr=ACCENT_RED, bold=True, sa=Pt(4))
para(tf1, "故障发生 → 告警触发 → 人工定位(30min+)\n→ 人工处置 → 事后复盘\nMTTR: 30分钟 ~ 数小时", sz=11, clr=DARK_TEXT, sa=Pt(0))

r2 = rect(sl, Inches(6.8), Inches(3.3), Inches(5.8), Inches(1.6), fill=RGBColor(0xF0,0xFD,0xF4))
tf2 = rich_tb(sl, Inches(6.95), Inches(3.35), Inches(5.5), Inches(1.5))
para(tf2, "目标模式 (主动预防)", sz=13, clr=ACCENT_GRN, bold=True, sa=Pt(4))
para(tf2, "AI预测风险 → 自动评估 → 预置防御\n→ 故障时自动自愈 → 自动复盘\nMTTR: <5分钟 (大部分<1分钟)", sz=11, clr=DARK_TEXT, sa=Pt(0))

# LLM横切
tb(sl, Inches(0.5), Inches(5.1), Inches(12), Inches(0.3),
   "LLM — 贯穿开发态与运行态的桥梁", sz=13, clr=PRIMARY, bold=True)
r = rect(sl, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.2))
rect_text(r,
    "开发态: 免测试故障定位(AgentFL, ICSE'24) → 自动程序修复(RepairAgent, ICSE'25) → 智能Fuzzing种子生成\n"
    "运行态: 日志分析(R-Log, arXiv'25, +228%) → 告警理解(Zha, Electronics'24) → 修复建议(Xpert, ICSE'24)\n"
    "策略: 高置信度自动执行 | 中置信度人机协同 | 低置信度人工决策",
    sz=10, clr=DARK_TEXT, align=PP_ALIGN.LEFT)

pg(sl, 10)

# ── 第11页：演进路线图 ───────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header(sl, "技术演进路线图", "短期(6月) → 中期(1-2年) → 长期(3-5年)")

phases = [
    ("短期 6个月", "立即可行", [
        "S1: 静态分析集成CI",
        "S2: ML异常检测(KPI)",
        "S3: 告警聚合机制",
        "S4: 协议Fuzzing(AFL++)",
        "S5: Beta因子基线评估",
    ], ACCENT_GRN),
    ("中期 1-2年", "重点突破", [
        "M1: Grid化同城多活",
        "M2: AI根因定位系统",
        "M3: LLM运维助手",
        "M4: 异构双栈部署",
        "M5: 智能话务预测",
    ], RGBColor(0xD9,0x77,0x06)),
    ("长期 3-5年", "架构变革", [
        "L1: 全国Grid化部署",
        "L2: 多云异构容灾",
        "L3: AI全自动自愈",
        "L4: 概率风险评估(PSA)",
        "L5: LLM持续可靠性工程",
    ], PRIMARY),
]
for i, (phase, tag, items, clr) in enumerate(phases):
    x = Inches(0.5+i*4.2)
    # 标题
    hdr = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.1), Inches(3.8), Inches(0.6))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = clr; hdr.line.fill.background()
    rect_text(hdr, phase, sz=16, clr=WHITE, bold=True)
    # 条目
    r = rect(sl, x, Inches(1.75), Inches(3.8), Inches(2.6))
    tf = rich_tb(sl, x+Inches(0.15), Inches(1.8), Inches(3.5), Inches(2.5))
    for j, item in enumerate(items):
        para(tf, f"● {item}", sz=11, clr=DARK_TEXT, sa=Pt(6))

# 量化目标表
tb(sl, Inches(0.5), Inches(4.6), Inches(12), Inches(0.3),
   "量化目标演进", sz=13, clr=PRIMARY, bold=True)
metrics = [
    ("指标", "当前", "短期", "中期", "长期"),
    ("可用性", "99.99%", "99.995%", "99.999%", "99.9995%"),
    ("MTTD", "5-15min", "<1min", "<30s", "<10s"),
    ("MTTR", "30-60min", "<10min", "<5min", "<1min"),
    ("Beta因子", "~0.3", "0.2", "0.05-0.15", "<0.01"),
    ("AI定位率", "<10%", "30%", "70%", ">95%"),
]
for i, row in enumerate(metrics):
    for j, cell in enumerate(row):
        x = Inches(0.5+j*2.4)
        y = Inches(5.0+i*0.33)
        b = True if i==0 or j==0 else False
        c = PRIMARY if i==0 else (MID_TEXT if j==0 else DARK_TEXT)
        tb(sl, x, y, Inches(2.3), Inches(0.28), cell, sz=9, clr=c, bold=b, align=PP_ALIGN.CENTER)

pg(sl, 11)

# ── 第12页：结论与建议 ───────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header(sl, "结论与行动建议", "面向云核心网产品决策者和架构师")

# 三个核心结论
conclusions = [
    ("方法论 × 工程实践", "核电纵深防御方法论 + 云Grid工程实践\n= 核心网\"超可靠架构蓝图\""),
    ("AI是倍增器非替代品", "先建立坚实架构基础(Grid/异构冗余)\n再叠加AI能力提升效率"),
    ("LLM贯穿开发与运行", "同时重塑软件开发(Bug定位/修复)\n和运行时(日志分析/自愈决策)"),
]
for i, (n, d) in enumerate(conclusions):
    x = Inches(0.5+i*4.2)
    r = rect(sl, x, Inches(1.1), Inches(3.8), Inches(1.5))
    tf = rich_tb(sl, x+Inches(0.15), Inches(1.15), Inches(3.5), Inches(1.4))
    para(tf, n, sz=14, clr=PRIMARY, bold=True, sa=Pt(6), align=PP_ALIGN.CENTER)
    para(tf, d, sz=11, clr=DARK_TEXT, sa=Pt(0), align=PP_ALIGN.CENTER)

# 五项建议
tb(sl, Inches(0.5), Inches(2.9), Inches(12), Inches(0.3),
   "五项行动建议", sz=16, clr=PRIMARY, bold=True)
actions = [
    ("1", "立即启动可靠性基线评估", "量化Beta因子、MTTD、MTTR → 为投入决策提供数据基础"),
    ("2", "6个月内完成\"低垂果实\"采摘", "静态分析集成 + AI异常检测 + 告警聚合 → 投入产出比最高"),
    ("3", "1年内启动Grid化同城多活", "中长期架构演进基础，先在非关键业务试点"),
    ("4", "建立跨领域可靠性工程团队", "核电方法论 + AI/ML + 云原生 + 电信核心网四领域"),
    ("5", "制定3-5年演进路线图并定期审视", "每季度审视，参考核电定期安全审查制度"),
]
for i, (num, title, desc) in enumerate(actions):
    y = Inches(3.35+i*0.65)
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), y, Inches(0.4), Inches(0.4))
    circ.fill.solid(); circ.fill.fore_color.rgb = PRIMARY; circ.line.fill.background()
    rect_text(circ, num, sz=14, clr=WHITE, bold=True)
    tb(sl, Inches(1.05), y, Inches(4.5), Inches(0.5), title, sz=13, clr=DARK_TEXT, bold=True)
    tb(sl, Inches(5.8), y, Inches(7), Inches(0.5), desc, sz=11, clr=MID_TEXT)

pg(sl, 12)

# ── 保存 ──────────────────────────────────────────────
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'output',
                   'cloud_core_reliability_insights_v2.pptx')
prs.save(out)
print(f"PPT已保存: {out}")
print(f"共 12 页")
